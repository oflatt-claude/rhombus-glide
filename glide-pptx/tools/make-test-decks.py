#!/usr/bin/env python3
"""Generate the .pptx fixtures used by the test suite.

Each deck isolates a slice of the format so a fidelity regression points at a
specific feature rather than at "slides look wrong".
"""
import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "tests", "decks")


def widescreen():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def blank(prs):
    """Layout 6 is 'Blank' in the default template."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def save(prs, name):
    path = os.path.abspath(os.path.join(OUT, name))
    prs.save(path)
    print("wrote", path)


# --------------------------------------------------------------------------
# 01: placeholders on the stock title/content layouts. Exercises inheritance
# of geometry, font and color from slideLayout -> slideMaster -> theme.
# --------------------------------------------------------------------------
def deck_placeholders():
    prs = widescreen()
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = "Deck Title"
    s.placeholders[1].text = "A subtitle inherited from the layout"

    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "Bulleted Content"
    tf = s.placeholders[1].text_frame
    tf.text = "First level bullet"
    for txt, lvl in [("Second level", 1), ("Third level", 2), ("Back to first", 0)]:
        p = tf.add_paragraph()
        p.text = txt
        p.level = lvl

    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Title Only"
    save(prs, "01-placeholders.pptx")


# --------------------------------------------------------------------------
# 02: explicit text boxes. All run/paragraph properties are on the slide, so
# this isolates text rendering from inheritance.
# --------------------------------------------------------------------------
def deck_text():
    prs = widescreen()
    s = blank(prs)

    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(6), Inches(3))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Plain 18pt"
    r.font.size = Pt(18)
    for label, mut in [
        ("Bold 24pt", lambda f: (setattr(f, "bold", True), setattr(f, "size", Pt(24)))),
        ("Italic 20pt", lambda f: (setattr(f, "italic", True), setattr(f, "size", Pt(20)))),
        ("Underline 16pt", lambda f: (setattr(f, "underline", True), setattr(f, "size", Pt(16)))),
        ("Red 28pt", lambda f: (setattr(f, "size", Pt(28)),
                                setattr(f.color, "rgb", RGBColor(0xC0, 0x20, 0x20)))),
        ("Courier New 18pt", lambda f: (setattr(f, "size", Pt(18)),
                                        setattr(f, "name", "Courier New"))),
    ]:
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = label
        mut(r.font)

    # One paragraph, several runs: tests run-level concatenation on a line.
    p = tf.add_paragraph()
    for txt, bold in [("mixed ", False), ("runs ", True), ("in ", False), ("one line", True)]:
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(18)
        r.font.bold = bold

    # Alignment: three boxes with a visible outline so drift is obvious.
    for i, algn in enumerate([PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.RIGHT]):
        tb = s.shapes.add_textbox(Inches(7), Inches(0.4 + i * 1.2), Inches(5.5), Inches(1.0))
        tb.text_frame.paragraphs[0].alignment = algn
        r = tb.text_frame.paragraphs[0].add_run()
        r.text = f"aligned {str(algn).split('.')[-1].split(' ')[0].lower()}"
        r.font.size = Pt(20)
        tb.line.color.rgb = RGBColor(0x99, 0x99, 0x99)
        tb.line.width = Pt(0.75)

    # Vertical anchoring inside a fixed box.
    for i, anch in enumerate([MSO_ANCHOR.TOP, MSO_ANCHOR.MIDDLE, MSO_ANCHOR.BOTTOM]):
        tb = s.shapes.add_textbox(Inches(0.5 + i * 2.2), Inches(4.2), Inches(2.0), Inches(2.6))
        tb.text_frame.vertical_anchor = anch
        r = tb.text_frame.paragraphs[0].add_run()
        r.text = str(anch).split(".")[-1].split(" ")[0].lower()
        r.font.size = Pt(16)
        tb.line.color.rgb = RGBColor(0x99, 0x99, 0x99)
        tb.line.width = Pt(0.75)

    save(prs, "02-text.pptx")


# --------------------------------------------------------------------------
# 03: autoshape geometry, fills and outlines.
# --------------------------------------------------------------------------
def deck_shapes():
    prs = widescreen()
    s = blank(prs)
    shapes = [
        MSO_SHAPE.RECTANGLE, MSO_SHAPE.ROUNDED_RECTANGLE, MSO_SHAPE.OVAL,
        MSO_SHAPE.ISOSCELES_TRIANGLE, MSO_SHAPE.RIGHT_TRIANGLE, MSO_SHAPE.DIAMOND,
        MSO_SHAPE.PENTAGON, MSO_SHAPE.HEXAGON, MSO_SHAPE.RIGHT_ARROW,
        MSO_SHAPE.LEFT_ARROW, MSO_SHAPE.CHEVRON, MSO_SHAPE.STAR_5_POINT,
    ]
    palette = [0x4472C4, 0xED7D31, 0xA5A5A5, 0xFFC000, 0x5B9BD5, 0x70AD47]
    for i, kind in enumerate(shapes):
        col, row = i % 4, i // 4
        sp = s.shapes.add_shape(
            kind, Inches(0.6 + col * 3.2), Inches(0.5 + row * 2.3),
            Inches(2.6), Inches(1.7))
        sp.fill.solid()
        sp.fill.fore_color.rgb = RGBColor.from_string(f"{palette[i % len(palette)]:06X}")
        sp.line.color.rgb = RGBColor(0x20, 0x20, 0x20)
        sp.line.width = Pt(1.5)
        sp.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        r = sp.text_frame.paragraphs[0].add_run()
        r.text = str(kind).split(".")[-1].split(" ")[0][:14]
        r.font.size = Pt(12)
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    s = blank(prs)
    # No fill / no line / thick line, and a rotated shape.
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.6), Inches(3), Inches(2))
    sp.fill.background()
    sp.line.color.rgb = RGBColor(0x00, 0x70, 0xC0)
    sp.line.width = Pt(4)

    sp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.2), Inches(0.6), Inches(3), Inches(2))
    sp.fill.solid()
    sp.fill.fore_color.rgb = RGBColor(0x70, 0xAD, 0x47)
    sp.line.fill.background()

    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.0), Inches(0.9), Inches(3), Inches(1.4))
    sp.fill.solid()
    sp.fill.fore_color.rgb = RGBColor(0xFF, 0xC0, 0x00)
    sp.line.color.rgb = RGBColor(0x00, 0x00, 0x00)
    sp.rotation = 20.0

    # Overlap order: later shapes paint on top.
    for i in range(4):
        sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(1.0 + i * 0.9), Inches(3.4), Inches(2.4), Inches(2.4))
        sp.fill.solid()
        sp.fill.fore_color.rgb = RGBColor.from_string(f"{palette[i]:06X}")
        sp.line.fill.background()

    conn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                  Inches(7.0), Inches(3.6), Inches(12.6), Inches(6.4))
    conn.line.color.rgb = RGBColor(0xC0, 0x00, 0x00)
    conn.line.width = Pt(3)
    conn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                  Inches(7.0), Inches(6.4), Inches(12.6), Inches(3.6))
    conn.line.color.rgb = RGBColor(0x00, 0x60, 0xC0)
    conn.line.width = Pt(1)
    save(prs, "03-shapes.pptx")


# --------------------------------------------------------------------------
# 04: pictures and groups.
# --------------------------------------------------------------------------
def deck_pictures(media_dir):
    prs = widescreen()
    s = blank(prs)
    s.shapes.add_picture(os.path.join(media_dir, "gradient.png"),
                         Inches(0.6), Inches(0.6), Inches(4), Inches(3))
    s.shapes.add_picture(os.path.join(media_dir, "checker.png"),
                         Inches(5.2), Inches(0.6), Inches(3), Inches(3))
    # Aspect change: same source stretched to a wide box.
    s.shapes.add_picture(os.path.join(media_dir, "checker.png"),
                         Inches(0.6), Inches(4.2), Inches(7.6), Inches(2.4))
    s.shapes.add_picture(os.path.join(media_dir, "gradient.png"),
                         Inches(9.0), Inches(1.0), Inches(3.6), Inches(2.7))

    s = blank(prs)
    # Group with a non-trivial child coordinate space (chOff/chExt != off/ext),
    # which is where naive group handling breaks.
    grp = s.shapes.add_group_shape()
    sp = grp.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1.2))
    sp.fill.solid()
    sp.fill.fore_color.rgb = RGBColor(0x44, 0x72, 0xC4)
    sp = grp.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.2), Inches(1.6), Inches(1.6), Inches(1.6))
    sp.fill.solid()
    sp.fill.fore_color.rgb = RGBColor(0xED, 0x7D, 0x31)
    tb = grp.shapes.add_textbox(Inches(1), Inches(2.6), Inches(3.6), Inches(0.8))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = "inside a group"
    r.font.size = Pt(18)
    # Scale the group: children must scale with it.
    grp.left, grp.top = Inches(6.5), Inches(2.0)
    grp.width, grp.height = Inches(6.0), Inches(4.0)
    save(prs, "04-pictures-groups.pptx")


# --------------------------------------------------------------------------
# 05: a plausible "real" deck -- the mix an actual presentation has.
# --------------------------------------------------------------------------
def deck_realistic(media_dir):
    prs = widescreen()
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = "Translating PowerPoint to Pict"
    s.placeholders[1].text = "A verified pptx -> Racket/Rhombus pipeline"

    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "Why bother?"
    tf = s.placeholders[1].text_frame
    tf.text = "PowerPoint is a good positioning tool"
    for txt in ["Code is a good animation tool",
                "Round-tripping keeps both",
                "PDF diffing keeps us honest"]:
        tf.add_paragraph().text = txt

    s = blank(prs)
    tb = s.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(11), Inches(1))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = "Pipeline"
    r.font.size = Pt(40)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0x1F, 0x3B, 0x63)

    stages = ["pptx", "IR", "pict", "PDF"]
    for i, name in enumerate(stages):
        sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.8 + i * 3.1), Inches(2.2), Inches(2.4), Inches(1.2))
        sp.fill.solid()
        sp.fill.fore_color.rgb = RGBColor(0x44, 0x72, 0xC4)
        sp.line.fill.background()
        sp.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        r = sp.text_frame.paragraphs[0].add_run()
        r.text = name
        r.font.size = Pt(22)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        if i < len(stages) - 1:
            sp = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                    Inches(3.3 + i * 3.1), Inches(2.55), Inches(0.5), Inches(0.5))
            sp.fill.solid()
            sp.fill.fore_color.rgb = RGBColor(0x88, 0x88, 0x88)
            sp.line.fill.background()

    s.shapes.add_picture(os.path.join(media_dir, "gradient.png"),
                         Inches(0.8), Inches(4.0), Inches(5.0), Inches(2.8))
    tb = s.shapes.add_textbox(Inches(6.4), Inches(4.0), Inches(6.2), Inches(2.8))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, txt in enumerate(["Shapes and text become picts.",
                             "Absolute positions become pin-over.",
                             "Everything else is a diff to explain."]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(20)
    save(prs, "05-realistic.pptx")


def make_media(media_dir):
    """Two small PNGs written without Pillow: one smooth, one hard-edged."""
    import struct
    import zlib

    def png(path, w, h, pixel):
        raw = bytearray()
        for y in range(h):
            raw.append(0)
            for x in range(w):
                raw.extend(pixel(x, y))
        def chunk(tag, data):
            return (struct.pack(">I", len(data)) + tag + data
                    + struct.pack(">I", zlib.crc32(tag + data) & 0xFFFFFFFF))
        with open(path, "wb") as f:
            f.write(b"\x89PNG\r\n\x1a\n")
            f.write(chunk(b"IHDR", struct.pack(">IIBBBBB", w, h, 8, 2, 0, 0, 0)))
            f.write(chunk(b"IDAT", zlib.compress(bytes(raw), 9)))
            f.write(chunk(b"IEND", b""))
        print("wrote", path)

    os.makedirs(media_dir, exist_ok=True)
    png(os.path.join(media_dir, "gradient.png"), 240, 180,
        lambda x, y: (int(255 * x / 239), int(255 * y / 179), 160))
    png(os.path.join(media_dir, "checker.png"), 160, 160,
        lambda x, y: ((30, 30, 40) if ((x // 20) + (y // 20)) % 2 else (240, 200, 60)))


def main():
    os.makedirs(OUT, exist_ok=True)
    media = os.path.abspath(os.path.join(OUT, "..", "media"))
    make_media(media)
    deck_placeholders()
    deck_text()
    deck_shapes()
    deck_pictures(media)
    deck_realistic(media)


if __name__ == "__main__":
    sys.exit(main())
